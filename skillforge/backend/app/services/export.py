import json
from io import BytesIO
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml


# ---------- Color Palette ----------
DARK_BG     = RGBColor(15, 23, 42)       # slate-900
BLUE_HEADER = RGBColor(30, 64, 175)      # blue-800
WHITE       = RGBColor(255, 255, 255)
LIGHT_TEXT  = RGBColor(248, 250, 252)    # slate-50
MUTED_TEXT  = RGBColor(148, 163, 184)    # slate-400
BODY_TEXT   = RGBColor(30, 41, 59)       # slate-800
GREEN       = RGBColor(22, 163, 74)       # green-600
CARD_BG     = RGBColor(241, 245, 249)    # slate-100
ACCENT      = RGBColor(59, 130, 246)     # blue-500


def _set_run_font(run, size_pt: int, bold: bool = False, color: RGBColor = BODY_TEXT):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_textbox(slide, left, top, width, height, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    return tf


def create_plan_pptx(plan_record) -> BytesIO:
    plan_data: Dict[str, Any] = json.loads(plan_record.plan_json)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # =========================================================
    # SLIDE 1 – Dark Title Slide
    # =========================================================
    slide = prs.slides.add_slide(blank)

    # Full dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Main title
    tf = _add_textbox(slide, 0.8, 2.1, 11.7, 1.4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = plan_data.get("title", "SkillForge Learning Plan")
    _set_run_font(run, 34, bold=True, color=LIGHT_TEXT)

    # Meta line
    tf = _add_textbox(slide, 0.8, 3.7, 11.7, 0.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        f"{plan_data.get('duration_weeks', '-')} Weeks   •   "
        f"{plan_data.get('total_hours', '-')} Hours   •   "
        f"Budget ₹{plan_data.get('recommended_budget_inr', 0)}"
    )
    _set_run_font(run, 18, color=MUTED_TEXT)

    # Footer
    tf = _add_textbox(slide, 0.8, 6.7, 11.7, 0.4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SkillForge  •  AI Career + Finance Co-Pilot for Engineering Students"
    _set_run_font(run, 12, color=MUTED_TEXT)

    # =========================================================
    # SLIDE 2 – Overview & Tips
    # =========================================================
    slide = prs.slides.add_slide(blank)

    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_HEADER
    header.line.fill.background()

    tf = _add_textbox(slide, 0.6, 0.28, 12, 0.55)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Plan Overview & Success Tips"
    _set_run_font(run, 24, bold=True, color=WHITE)

    # Tips as clean list
    tips = plan_data.get("tips", [])
    tf = _add_textbox(slide, 0.8, 1.5, 11.7, 5.3)

    for i, tip in enumerate(tips):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(16)
        run = p.add_run()
        run.text = f"→   {tip}"
        _set_run_font(run, 17, color=BODY_TEXT)

    # =========================================================
    # WEEKLY SLIDES
    # =========================================================
    weekly = plan_data.get("weekly_plan", [])

    for week in weekly:
        slide = prs.slides.add_slide(blank)

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
        header.fill.solid()
        header.fill.fore_color.rgb = BLUE_HEADER
        header.line.fill.background()

        tf = _add_textbox(slide, 0.6, 0.28, 12, 0.55)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Week {week.get('week', '-')}   •   Focus: {week.get('focus_skill', '')}"
        _set_run_font(run, 22, bold=True, color=WHITE)

        # ---- Left Card: Tasks ----
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        tf = _add_textbox(slide, 0.7, 1.55, 5.6, 0.4)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "TASKS"
        _set_run_font(run, 14, bold=True, color=BLUE_HEADER)

        tf = _add_textbox(slide, 0.7, 2.1, 5.6, 4.5)
        for i, task in enumerate(week.get("tasks", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = f"•  {task}"
            _set_run_font(run, 14, color=BODY_TEXT)

        # ---- Right Card: Resources + Meta ----
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        tf = _add_textbox(slide, 7.0, 1.55, 5.6, 0.4)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "RESOURCES"
        _set_run_font(run, 14, bold=True, color=BLUE_HEADER)

        tf = _add_textbox(slide, 7.0, 2.1, 5.6, 2.8)
        for i, res in enumerate(week.get("resources", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = f"•  {res}"
            _set_run_font(run, 14, color=BODY_TEXT)

        # Cost + Hours at bottom of right card
        tf = _add_textbox(slide, 7.0, 5.3, 5.6, 1.2)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Estimated Cost:  ₹{week.get('estimated_cost_inr', 0)}"
        _set_run_font(run, 15, bold=True, color=GREEN)

        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"Daily Hours:  {week.get('daily_hours', 2)} hrs"
        _set_run_font(run, 14, color=BODY_TEXT)

    # =========================================================
    # FINAL SLIDE
    # =========================================================
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tf = _add_textbox(slide, 1.0, 2.5, 11.3, 2.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Stay consistent.\nTrack your spending.\nShip projects every week."
    _set_run_font(run, 28, bold=True, color=LIGHT_TEXT)

    tf = _add_textbox(slide, 1.0, 5.2, 11.3, 0.5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SkillForge  •  Built for Engineering Students"
    _set_run_font(run, 16, color=MUTED_TEXT)

    # ===== Always return the buffer =====
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer